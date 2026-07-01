"""
Document text extraction tools.
Extracts text from: PDF, DOCX, XLSX, PPTX, images (OCR), text/code files.
"""
import os
import subprocess
import tempfile
import re
from datetime import datetime
from typing import Dict, Optional
from tools.base import Tool, ToolResult


def _format_size(size_bytes: int) -> str:
    if size_bytes == 0:
        return "0 B"
    for unit in ["B", "KB", "MB", "GB", "TB"]:
        if size_bytes < 1024:
            return f"{size_bytes:.1f} {unit}" if unit != "B" else f"{size_bytes} B"
        size_bytes /= 1024
    return f"{size_bytes:.1f} PB"


# -- Detect available backends --
HAS_PDFTOTEXT = subprocess.run(["which", "pdftotext"], capture_output=True, text=True).returncode == 0

HAS_TESSERACT = subprocess.run(["which", "tesseract"], capture_output=True, text=True).returncode == 0

HAS_DOCX = False
HAS_XLSX = False
HAS_PPTX = False
HAS_PYPDF = False
HAS_PDFMINER = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    pass

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    pass

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    pass

try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    pass

try:
    from pdfminer.high_level import extract_text as pdfminer_extract
    HAS_PDFMINER = True
except ImportError:
    pass


SUPPORTED_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".css",
    ".json", ".yaml", ".yml", ".xml", ".toml", ".ini", ".cfg", ".conf",
    ".sh", ".bash", ".zsh", ".fish", ".bat", ".ps1",
    ".c", ".cpp", ".h", ".hpp", ".java", ".go", ".rs", ".rb", ".php",
    ".swift", ".kt", ".scala", ".clj", ".ex", ".exs",
    ".sql", ".r", ".m", ".mm", ".pl", ".pm", ".lua", ".hs",
    ".csv", ".log",
}

BINARY_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp"}


def extract_text_from_file(file_path: str, max_pages: int = 20, ocr_language: str = "eng") -> dict:
    """Extract text from a file based on its extension.
    
    Returns dict with keys: text (str), method (str), pages (int), error (str, optional)
    """
    ext = os.path.splitext(file_path)[1].lower()
    result = {"text": "", "method": "unknown", "pages": 0, "error": None}

    # Plain text and code files
    if ext in SUPPORTED_EXTENSIONS:
        try:
            with open(file_path, "r", errors="replace") as f:
                content = f.read()
            result["text"] = content
            result["method"] = "plain_text"
            result["pages"] = content.count("\n") + 1
            return result
        except Exception as e:
            result["error"] = str(e)
            return result

    # PDF
    if ext == ".pdf":
        return _extract_pdf(file_path, max_pages)

    # Office documents
    if ext == ".docx":
        return _extract_docx(file_path)
    if ext == ".xlsx":
        return _extract_xlsx(file_path)
    if ext == ".pptx":
        return _extract_pptx(file_path)

    # Images (OCR)
    if ext in IMAGE_EXTENSIONS:
        return _extract_image_ocr(file_path, ocr_language)

    result["error"] = f"Unsupported file type: {ext}"
    return result


def _extract_pdf(file_path: str, max_pages: int = 20) -> dict:
    result = {"text": "", "method": "pdf", "pages": 0, "error": None}

    # Try pdftotext first (fastest, most reliable)
    if HAS_PDFTOTEXT:
        try:
            cmd = ["pdftotext", "-f", "1", "-l", str(max_pages), file_path, "-"]
            proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            if proc.returncode == 0:
                text = proc.stdout
                # Count pages from output markers
                pages = len(re.findall(r'\f', text)) + 1 if text else 0
                result["text"] = text.strip()
                result["pages"] = pages or 1
                return result
            result["error"] = proc.stderr.strip()
        except subprocess.TimeoutExpired:
            result["error"] = "pdftotext timed out"
        except Exception as e:
            result["error"] = str(e)

    # Fallback: pypdf
    if HAS_PYPDF:
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            pages = min(len(reader.pages), max_pages)
            texts = []
            for i in range(pages):
                texts.append(reader.pages[i].extract_text() or "")
            result["text"] = "\n\n".join(texts).strip()
            result["pages"] = pages
            result["method"] = "pypdf"
            return result
        except Exception as e:
            result["error"] = result.get("error") or str(e)

    # Fallback: pdfminer
    if HAS_PDFMINER:
        try:
            text = pdfminer_extract(file_path, maxpages=max_pages)
            result["text"] = text.strip()
            result["pages"] = max_pages
            result["method"] = "pdfminer"
            return result
        except Exception as e:
            result["error"] = result.get("error") or str(e)

    if not result["error"]:
        result["error"] = "No PDF extraction backend available (install pdftotext, pypdf, or pdfminer.six)"
    return result


def _extract_docx(file_path: str) -> dict:
    result = {"text": "", "method": "docx", "pages": 0, "error": None}
    if not HAS_DOCX:
        result["error"] = "python-docx not installed"
        return result
    try:
        import docx
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        # Also extract tables
        table_texts = []
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                table_texts.append(" | ".join(cells))
        all_text = "\n".join(paragraphs)
        if table_texts:
            all_text += "\n\n--- TABLES ---\n" + "\n".join(table_texts)
        result["text"] = all_text.strip()
        result["pages"] = len(paragraphs) // 40 + 1
        return result
    except Exception as e:
        result["error"] = str(e)
        return result


def _extract_xlsx(file_path: str) -> dict:
    result = {"text": "", "method": "xlsx", "pages": 0, "error": None}
    if not HAS_XLSX:
        result["error"] = "openpyxl not installed"
        return result
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            sheet_content = "\n".join(rows)
            sheets_text.append(f"--- Sheet: {sheet_name} ---\n{sheet_content}")
        result["text"] = "\n\n".join(sheets_text).strip()
        result["pages"] = len(wb.sheetnames)
        wb.close()
        return result
    except Exception as e:
        result["error"] = str(e)
        return result


def _extract_pptx(file_path: str) -> dict:
    result = {"text": "", "method": "pptx", "pages": 0, "error": None}
    if not HAS_PPTX:
        result["error"] = "python-pptx not installed"
        return result
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_content.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_content.append(" | ".join(cells))
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_content))
        result["text"] = "\n\n".join(slides_text).strip()
        result["pages"] = len(prs.slides)
        return result
    except Exception as e:
        result["error"] = str(e)
        return result


def _extract_image_ocr(file_path: str, language: str = "eng") -> dict:
    result = {"text": "", "method": "ocr", "pages": 1, "error": None}
    if not HAS_TESSERACT:
        result["error"] = "tesseract OCR not available (install tesseract-ocr and pytesseract)"
        return result
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang=language)
        result["text"] = text.strip()
        return result
    except Exception as e:
        result["error"] = str(e)
        return result


class ExtractDocumentTextTool(Tool):
    name = "extract_document_text"
    description = (
        "Extract text content from documents: PDF, DOCX, XLSX, PPTX, "
        "images (OCR), and plain text/code files. Returns the full text."
    )
    parameters_schema = {
        "path": "Absolute path to the document file",
        "max_pages": "(optional) Max pages to extract from PDFs. Default 20.",
        "ocr_language": "(optional) Tesseract OCR language. Default 'eng'.",
    }

    async def execute(self, args: Dict) -> ToolResult:
        path = os.path.expanduser(args.get("path", ""))
        max_pages = int(args.get("max_pages", 20))
        ocr_language = args.get("ocr_language", "eng")

        if not os.path.exists(path):
            return ToolResult(success=False, message=f"File does not exist: {path}")
        if os.path.isdir(path):
            return ToolResult(success=False, message=f"Not a file: {path}")

        ext = os.path.splitext(path)[1].lower()
        st = os.stat(path)
        extraction = extract_text_from_file(path, max_pages, ocr_language)

        if extraction["error"] and not extraction["text"]:
            return ToolResult(
                success=False,
                message=f"Extraction failed: {extraction['error']}",
                data={
                    "path": path,
                    "extension": ext,
                    "size": st.st_size,
                    "size_formatted": _format_size(st.st_size),
                },
            )

        return ToolResult(
            success=True,
            data={
                "path": path,
                "name": os.path.basename(path),
                "extension": ext,
                "size": st.st_size,
                "size_formatted": _format_size(st.st_size),
                "method": extraction["method"],
                "pages": extraction["pages"],
                "text": extraction["text"],
                "text_length": len(extraction["text"]),
                "modified": datetime.fromtimestamp(st.st_mtime).isoformat(),
            },
            message=f"Extracted {len(extraction['text'])} chars from {os.path.basename(path)} ({extraction['method']})",
            files_affected=[path],
        )


class BatchExtractTextTool(Tool):
    name = "batch_extract_text"
    description = (
        "Find and extract text from all supported documents in a directory. "
        "Returns extracted text from each file."
    )
    parameters_schema = {
        "path": "Directory to scan",
        "recursive": "(optional) Search recursively. Default true.",
        "extensions": "(optional) Comma-separated extensions to include. Default all supported.",
        "max_files": "(optional) Max files to process. Default 50.",
        "max_pages": "(optional) Max pages per PDF. Default 20.",
        "ocr_language": "(optional) Tesseract OCR language. Default 'eng'.",
    }

    async def execute(self, args: Dict) -> ToolResult:
        path = os.path.expanduser(args.get("path", ""))
        recursive = args.get("recursive", True)
        extensions_str = args.get("extensions", "")
        max_files = int(args.get("max_files", 50))
        max_pages = int(args.get("max_pages", 20))
        ocr_language = args.get("ocr_language", "eng")

        if not os.path.exists(path):
            return ToolResult(success=False, message=f"Path does not exist: {path}")
        if not os.path.isdir(path):
            return ToolResult(success=False, message=f"Not a directory: {path}")

        # Determine which extensions to include
        ALL_EXTRACTABLE = SUPPORTED_EXTENSIONS | BINARY_DOCUMENT_EXTENSIONS | IMAGE_EXTENSIONS
        if extensions_str:
            target_exts = set()
            for ext in extensions_str.split(","):
                ext = ext.strip().lower()
                if not ext.startswith("."):
                    ext = "." + ext
                target_exts.add(ext)
        else:
            target_exts = ALL_EXTRACTABLE

        # Walk directory
        files_found = []
        try:
            if recursive:
                for root, dirs, fnames in os.walk(path):
                    for fname in fnames:
                        if len(files_found) >= max_files:
                            break
                        ext = os.path.splitext(fname)[1].lower()
                        if ext in target_exts:
                            files_found.append(os.path.join(root, fname))
                    if len(files_found) >= max_files:
                        break
            else:
                for fname in os.listdir(path):
                    if len(files_found) >= max_files:
                        break
                    full = os.path.join(path, fname)
                    if os.path.isfile(full):
                        ext = os.path.splitext(fname)[1].lower()
                        if ext in target_exts:
                            files_found.append(full)
        except PermissionError:
            return ToolResult(success=False, message=f"Permission denied: {path}")

        if not files_found:
            return ToolResult(
                success=True,
                data={"results": [], "total": 0, "search_path": path},
                message="No supported documents found",
            )

        # Extract text from each file
        results = []
        errors = []
        for fp in files_found:
            try:
                extraction = extract_text_from_file(fp, max_pages, ocr_language)
                st = os.stat(fp)
                results.append({
                    "path": fp,
                    "name": os.path.basename(fp),
                    "extension": os.path.splitext(fp)[1].lower(),
                    "size": st.st_size,
                    "size_formatted": _format_size(st.st_size),
                    "method": extraction.get("method", "unknown"),
                    "text_length": len(extraction.get("text", "")),
                    "text_preview": extraction.get("text", "")[:500],
                    "error": extraction.get("error"),
                })
                if extraction.get("error"):
                    errors.append({"path": fp, "error": extraction["error"]})
            except Exception as e:
                errors.append({"path": fp, "error": str(e)})

        return ToolResult(
            success=True,
            data={
                "results": results,
                "total": len(results),
                "errors": errors,
                "search_path": path,
            },
            message=f"Extracted text from {len(results)} files ({len(errors)} errors)",
            files_affected=[path],
        )


ALL_EXTRACTION_TOOLS = [
    ExtractDocumentTextTool(),
    BatchExtractTextTool(),
]
